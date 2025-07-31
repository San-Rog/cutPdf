import pymupdf
import streamlit as st
import streamlit.components.v1 as components
import zipfile
import os
import time
import textwrap
import xlsxwriter
import numpy as np
import pandas as pd
import random
from segno import helpers
import subprocess
import datetime
from PyPDF2 import PdfReader, PdfWriter
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Pt
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
    
@st.cache_data   
def nameFile():
    symbols = ['-', ':', '.']
    nowTime = str(datetime.datetime.now())
    try:
        for symbol in symbols: 
            nowTime = nowTime.replace(symbol, '_')
    except:
        pass
    return nowTime
    
@st.cache_data  
def extractText(filePdf):
    text = ''
    docPdf = pymupdf.open(filePdf)
    for page in docPdf:
       text += page.get_text()
    docPdf.close()
    return text
        
@st.cache_data
def extractUrls(filePdf):
    docPdf = pymupdf.open(filePdf)
    allLinks = []
    for p, page in enumerate(docPdf):
        links = page.get_links()
        for link in links:
            nameUrl = link["uri"]
            fromUrl = link["from"]
            newText = f'{nameUrl}; coordenadas: {fromUrl}\n'
            allLinks.append(newText)
    text = ''.join(allLinks) 
    docPdf.close()
    return text
    
def mensResult(value, nFiles, modelButt, fileTmp, fileFinal):
    opt = st.session_state[listKeys[5]]
    if opt == 0:
        crt = optionsSel[3]
    else:
        crt = optionsSel[opt]    
    colMens, colDown = st.columns([8, 2]) 
    if value == 1:
        if modelButt == 'zip': 
                with open(fileTmp, "rb") as file:
                    colDown.download_button(label="Download",
                                            data=file,
                                            file_name=fileFinal,
                                            mime='application/zip', 
                                            icon=":material/download:", 
                                            use_container_width=True)
        colMens.success(f'Gerado o zipado :blue[**{fileFinal}**] com ***{nFiles}*** arquivo(s) (:red[**{crt}**]). Clique no botão ao lado 👉.', 
                        icon='✔️') 
    elif value == 0:
        colDown.download_button(label='Download', 
                                data=fileTmp,
                                file_name=fileFinal,
                                mime='application/octet-stream', 
                                icon=":material/download:", 
                                use_container_width=True)
        colMens.success(f'Gerado o arquivo :blue[**{fileFinal}**] (:red[**{crt}**]). Clique no botão ao lado 👉.', 
                        icon='✔️') 
    elif value == 2:
        colDown.download_button(label='Download',
                                data=fileTmp,
                                file_name=fileFinal,
                                mime="text/csv", 
                                icon=":material/download:", 
                                use_container_width=True)
        colMens.success(f'Gerado o arquivo :blue[**{fileFinal}**] (:red[**{crt}**]). Clique no botão ao lado 👉.', 
                        icon='✔️')
    elif value == 3:
        colDown.download_button(label='Download',
                                data=fileTmp,
                                file_name=fileFinal,
                                mime='application/octet-stream', 
                                use_container_width=True)
        colMens.success(f'Gerado o arquivo :blue[**{fileFinal}**] (:red[**{crt}**]). Clique no botão ao lado 👉.', 
                        icon='✔️')
    elif value == 4:
        colDown.download_button(label='Download',
                                data=fileTmp,
                                file_name=fileFinal,
                                 mime='application/vnd.openxmlformats-officedocument.wordprocessingml.document', 
                                use_container_width=True)
        colMens.success(f'Gerado o arquivo :blue[**{fileFinal}**] (:red[**{crt}**]). Clique no botão ao lado 👉.', 
                        icon='✔️')
    
def extractImgs(filePdf):
    docPdf = pymupdf.open(filePdf)
    allImgName = []
    for p, page in enumerate(docPdf):
        imageList = page.get_images(full=True)
        if imageList:
            for i, img in enumerate(imageList):
                xref = img[0]
                baseImg = docPdf.extract_image(xref)
                imgBytes = baseImg["image"]
                imgExt = baseImg["ext"]
                imgName = f"image_{p+1}_{i+1}.{imgExt}"
                with open(imgName, "wb") as fileImg:
                    fileImg.write(imgBytes)
                allImgName.append(imgName)
    return allImgName

def downloadExt(files, namePdf, numPgOne, numPgTwo, obj):
    fileTmp = f'{nameFile()}_tempFile.zip'
    fileZip = f'{namePdf}_{numPgOne}_{numPgTwo}_{nameFile()}.zip'
    for file in files:
        with open(file, "rb") as extFile:
           PDFbyte = extFile.read()
        with zipfile.ZipFile(fileTmp, 'a') as extFile:
           extFile.writestr(file, PDFbyte)
    nFiles = len(files) 
    if nFiles > 0:
        mensResult(1, len(files), 'zip', fileTmp, fileZip)
    else:
        strEmpty = f'😢 Extração fracassada!\n🔴 arquivo {namePdf} \nsem {obj} extraível no intervalo de páginas {numPgOne}-{numPgTwo}!'
        config(strEmpty)

def rotatePdf(filePdf, index):
    inputPdf = filePdf
    name, ext = os.path.splitext(inputPdf)
    angle = int(valAngles[index].replace('°', ''))
    output = f'{name}_rotate_{angle}{ext}'
    docPdf = pymupdf.open(filePdf)
    for page in docPdf:
        page.set_rotation(angle)
    docPdf.save(output)
    docPdf.close()
    return output
    
def saveAllPdf(outputBase, contPartes, writer):
    outputPdf = f"{outputBase}{contPartes + 1}.pdf"
    with open(outputPdf, "wb") as outputFile:
        writer.write(outputFile)
    docPdf = pymupdf.open(outputFile)
    countPg.append(len(docPdf))
    docPdf.close()
    return outputPdf

def divideBySize(inputPdf, sizeMax, outputBase):
    filesCutSave = []
    try:
        reader = PdfReader(inputPdf)
        nPgs = len(reader.pages)
        sizeActual = 0
        contPartes = 0
        writer = PdfWriter()
        for i in range(nPgs):
            nameTeste = f'teste_{i+1}.pdf'
            namesTeste.append(nameTeste)
            page = reader.pages[i]
            writer.add_page(page) 
            with open(nameTeste, 'wb') as g:
                writer.write(g)
            sizeActual = os.path.getsize(nameTeste)/(1024**2)
            if sizeActual >= sizeMax:
                outputPdf = saveAllPdf(outputBase, contPartes, writer)
                filesCutSave.append(outputPdf)
                writer = PdfWriter()
                sizeActual = 0
                contPartes += 1
        if len(writer.pages) > 0:
            outputPdf = saveAllPdf(outputBase, contPartes, writer)
            filesCutSave.append(outputPdf)
    except Exception as e:
        st.error(f"Ocorreu um erro: {e} - página {i+1}", icon='🛑')
    return filesCutSave    

def createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, rotate):
    numPgOne -= 1    
    inputPdf = docPdf
    name, ext = os.path.splitext(namePdf)
    outputPdf = f'{name}_{numPgOne + 1}_{numPgTwo}.pdf'
    listSel = [pg for pg in range(numPgOne, numPgTwo)]
    docPdf.select(listSel)
    docPdf.save(outputPdf)
    if rotate:
        outputPdf = rotatePdf(outputPdf, index) 
    docPdf.close()
    return outputPdf   

def addWatermark(inputPdf, valMark):
    doc = pymupdf.open(inputPdf)
    for page_num in range(doc.page_count):
        page = doc[page_num]
        page_rect = page.rect
        x = page_rect.width/6.5
        y = page_rect.height * 0.95
        page.insert_text(
            (x, y),  
            valMark,  
            fontsize=25,
            color=(0.7, 0.7, 0.4),  
            rotate=0,  
            fill_opacity=0.3,
            stroke_opacity=0.3
        )
    doc.save(inputPdf, incremental=True, encryption=0)
    doc.close()
    return inputPdf

def selPdfMark(docPdf, numPgOne, numPgTwo, namePdf, index, valMark):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, True)
    pdfMark = addWatermark(outputPdf, valMark)
    downPdfUnique(pdfMark, numPgOne, numPgTwo, namePdf)           
    
def selPgsSize(docPdf, numPgOne, numPgTwo, namePdf, index, sizeMax):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, True)
    inputPdf = outputPdf
    sizeMaxStr = str(sizeMax).replace('.', '_')
    sizeSplit = sizeMaxStr.split('_')
    try:
        numOne = sizeSplit[0]
        numTwo = sizeSplit[1][:2]
        if numTwo.strip() == '00':
             numTwo = ''
        sizeMaxStr = numOne + '_' + numTwo 
    except:
        pass
    outputBase = f'{os.path.splitext(inputPdf)[0]}_divisão_{sizeMaxStr}_Mb__parte_'
    filesCutSave = divideBySize(inputPdf, sizeMax, outputBase)
    downloadExt(filesCutSave, namePdf, numPgOne, numPgTwo, 'pedaços')

@st.cache_data
def extractTables(filePdf):
    docPdf = pymupdf.open(filePdf)
    AllTable = []
    for page in docPdf:
        tabs = page.find_tables()
        for t, tab in enumerate(tabs):
            listaTable = tab.extract()
            for lista in listaTable:
                AllTable.append(lista)
    return AllTable
    
def selImgUrlsPgs(docPdf, numPgOne, numPgTwo, namePdf, mode, index):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, False)
    filesImg = extractImgs(outputPdf)
    downloadExt(filesImg, namePdf, numPgOne, numPgTwo, 'imagens')
    
def selTablesPgs(docPdf, numPgOne, numPgTwo, namePdf, index):
    name, ext = os.path.splitext(namePdf)
    newName = f'{name}_{numPgOne}_{numPgTwo}.xlsx'
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, False)
    allTables = extractTables(outputPdf)
    fileFinal = newName
    workbook = xlsxwriter.Workbook(fileFinal)
    worksheet = workbook.add_worksheet('aba_dados')
    for rowNum, rowData in enumerate(allTables):
        worksheet.write_row(rowNum, 0, rowData)
    workbook.close()
    with open(fileFinal, "rb") as file:
        byFinal = file.read()
    mensResult(3, 0, 'xlsx', byFinal, fileFinal)

@st.cache_data   
def imagesConvert(filePdf):
    docPdf = pymupdf.open(filePdf)
    nPages = len(docPdf)
    listImgs = []
    for pg in range(nPages):
        page = docPdf.load_page(pg)
        pix = page.get_pixmap()
        fileImg = f'imagem_{pg + 1}.png'
        pix.save(fileImg)
        listImgs.append(fileImg)
    return listImgs  

@st.cache_data 
def docxConvert(filePdf):
    name = os.path.splitext(filePdf)[0]
    fileDocx = f'{name}.docx'
    try:
        cv = Converter(filePdf)
        cv.convert(fileDocx, start=0, end=None)
        cv.close()
    except: 
        pass
    return fileDocx
 
@st.cache_data 
def ppTxConvert(filePdf):
    docPdf = pymupdf.open(filePdf)
    baseName = os.path.basename(filePdf)
    name, ext = os.path.splitext(baseName)
    newName = f'{name}.pptx'
    dictAllTexts = {}
    for pg, page in enumerate(docPdf):
        nPg = pg + 1
        text = page.get_text()
        dictAllTexts.setdefault(nPg, '')
        dictAllTexts[nPg] += f'{text}\n'
    wrapper = textwrap.TextWrapper(width=75)
    p = Presentation()
    contPg = 0
    for dctAll, texts in dictAllTexts.items():
        textSplit = [txt.strip() for txt in texts.split('\n') if len(txt.strip()) != 0]
        textAdd = ''
        contSeg = 0 
        for tx, text in enumerate(textSplit):
            if tx%14 == 0 and tx != 0:
                contSeg += 1
                s = p.slides.add_slide(p.slide_layouts[5])
                titlePara = s.shapes.title.text_frame.paragraphs[0]
                titlePara.font.name = "Times New Roman"
                titlePara.font.size = Pt(18)
                titlePara.text = f'Arquivo {name} - página {dctAll} - segmento {contSeg}'
                txt_box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1))
                txt_frame = txt_box.text_frame
                n = txt_frame.add_paragraph()
                string = wrapper.fill(text=textAdd)
                n.text = string
                n.alignment = PP_ALIGN.JUSTIFY
                textAdd = ''
                textAdd += text + '\n'
            else:
                textAdd += text + '\n'
        s = p.slides.add_slide(p.slide_layouts[5])
        titlePara = s.shapes.title.text_frame.paragraphs[0]
        titlePara.font.name = "Times New Roman"
        titlePara.font.size = Pt(18)
        titlePara.text = f'Arquivo {name} - página {dctAll} - segmento {contSeg}'
        txt_box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1))
        txt_frame = txt_box.text_frame
        n = txt_frame.add_paragraph()
        string = wrapper.fill(text=textAdd)
        n.text = string
        n.alignment = PP_ALIGN.JUSTIFY
    p.save(newName)
    return newName

@st.cache_data   
def createImgQrCode():
    fileImg = 'myContact.png'
    valuesQrcode = []
    for k, key in enumerate(qrCodeKeys):
        valueState = st.session_state[key]
        if type(valueState) == tuple:
            valueState = valueState[0]
        if len(valueState.strip()) == 0:
            valuesQrcode.append(valuesReserve[k])
        else:
            valuesQrcode.append(valueState)    
    qrcode = helpers.make_mecard(name=valuesQrcode[0], 
                                 phone=valuesQrcode[1], 
                                 email=valuesQrcode[2])
    qrcode.save(fileImg, scale=1)
    return fileImg    
    
def selPdfToImg(docPdf, numPgOne, numPgTwo, namePdf, index): 
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, True)
    listImgs = imagesConvert(outputPdf)
    downloadExt(listImgs, namePdf, numPgOne, numPgTwo, 'pdf_img')
    
def selPdfToDocx(docPdf, numPgOne, numPgTwo, namePdf, index): 
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, False)
    fileDocx = docxConvert(outputPdf)
    with open(fileDocx, "rb") as file:
        byFinal = file.read()
    mensResult(4, 0, 'docx', byFinal, fileDocx)
    
def selPdfToPPtx(docPdf, numPgOne, numPgTwo, namePdf, index):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, False)
    filePptx = ppTxConvert(outputPdf)
    with open(filePptx, "rb") as file:
        byFinal = file.read()
    mensResult(4, 0, 'pptx', byFinal, filePptx)
    
def selPdfToQrcode(docPdf, numPgOne, numPgTwo, namePdf, index):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, True)
    fileImg = createImgQrCode()
    filePdf = insertImgPdf(outputPdf, fileImg)
    with open(filePdf, "rb") as file:
        PDFbyte = file.read()
    mensResult(0, 1, 'pdf', PDFbyte, filePdf)    

def insertImgPdf(filePdf, imgFile):
    baseName = os.path.basename(filePdf)
    name, ext = os.path.splitext(baseName)
    newPdf = f'{name}_img.pdf'
    headY = 720
    docPdf = pymupdf.open(filePdf)
    for pg in range(len(docPdf)):
        page = docPdf.load_page(pg)
        rect = page.rect  
        img = pymupdf.open(imgFile)
        img_rect = img.load_page(0).rect
        x0 = (rect.width - img_rect.width) / 13  
        y0 = rect.height - headY / 11            
        page.insert_image((x0, y0, x0 + img_rect.width, y0 + img_rect.height), filename=imgFile)             
    docPdf.save(newPdf)
    return newPdf    
    
def selTxtUrlPgs(docPdf, numPgOne, numPgTwo, namePdf, mode, index):
    outputPdf = createPdfSel(docPdf, numPgOne, numPgTwo, namePdf, index, False)
    if mode == 0:
        text = extractText(outputPdf)
        strLabel = "Download_text"
        outputTxt = f'{namePdf}_{numPgOne}_{numPgTwo}_text.txt'
        strEmpty = f'😢 Extração fracassada!\n🔴 arquivo {namePdf} \nsem texto extraível no intervalo de páginas {numPgOne}-{numPgTwo}!'
    else:
        text = extractUrls(outputPdf)
        strLabel = "Download_urls"
        outputTxt = f'{namePdf}_{numPgOne}_{numPgTwo}_urls.txt'
        strEmpty = f'😢 Extração fracassada!\n🔴 arquivo {namePdf} \nsem URL extraível no intervalo de páginas {numPgOne}-{numPgTwo}!'
    if len(text.strip()) > 0:
        mensResult(2, 1, 'txt', text, outputTxt)        
    else:
        config(strEmpty)    
                                          
def selDelPgs(docPdf, numPgOne, numPgTwo, namePdf, mode, index):
    numPgOne -= 1
    inputPdf = docPdf
    name, ext = os.path.splitext(namePdf)
    listPgs = seqPages(numPgOne, numPgTwo)
    if mode == 0:
        outputPdf = f'{name}_sel_{numPgOne + 1}_{numPgTwo}{ext}'
        listSel = [pg for pg in range(numPgOne, numPgTwo) if pg in listPgs]
    else:
        numPages = inputPdf.page_count
        outputPdf = f'{name}_del_{numPgOne + 1}_{numPgTwo}{ext}'
        listSel = [pg for pg in range(numPages) if pg not in range(numPgOne, numPgTwo)]
        listPlus = [pg for pg in range(numPgOne, numPgTwo) if pg not in listPgs]
        listSel = listPlus + listSel
    docPdf.select(listSel)
    docPdf.save(outputPdf)
    docPdf.close()
    if index != 4:
        outputPdf = rotatePdf(outputPdf, index)       
    downPdfUnique(outputPdf, numPgOne, numPgTwo, namePdf) 
                       
def downPdfUnique(outputPdf, numPgOne, numPgTwo, namePdf):
    with open(outputPdf, "rb") as pdf_file:
        PDFbyte = pdf_file.read()
    if len(PDFbyte) > 0:
        mensResult(0, 1, 'pdf', PDFbyte, outputPdf)
        
def extractPgs(docPdf, numPgOne, numPgTwo, mode, namePdf, index):
    numPgOne -= 1
    filesPdf = [docPdf]
    filesRead = []  
    name, ext = os.path.splitext(namePdf)
    for file in filesPdf:
        numPages = file.page_count
        diffPg = abs(numPgTwo - numPgOne)
        minPg = min([numPgOne, numPgTwo])
        listPg = [pg for pg in range(minPg, diffPg)]
        for p, pageNum in enumerate(listPg):
            inputPdf = file
            outputPdf = f'{name}_{pageNum + 1}.pdf'
            newPdf = pymupdf.open()
            newPdf.insert_pdf(inputPdf, from_page=pageNum, to_page=pageNum)
            newPdf.save(outputPdf)
            if index != 4:
                outputPdf = rotatePdf(outputPdf, index) 
            filesRead.append(outputPdf)
            newPdf.close()
    downloadExt(filesRead, namePdf, numPgOne, numPgTwo, 'páginas')

def exibeInfo(docPdf):
    @st.dialog(' ')
    def config():
        trace= '_'*10
        nPgs = docPdf.page_count
        size = round(uploadPdf.size/1024, 2)
        if size > 1024:
            size /= 1024
            size = round(size, 2)
            unit = 'MB'
        else:
            unit = 'KB'
        typFile = uploadPdf.type
        dirty = docPdf.is_dirty
        pdfYes = docPdf.is_pdf
        close = docPdf.is_closed
        formPdf = docPdf.is_form_pdf
        encry = docPdf.is_encrypted
        pdfMeta = docPdf.metadata
        st.markdown(f'🗄️ **Tamanho**: {size}{unit}') 
        st.markdown(f'📄️ **Total de páginas**: {nPgs}')
        dictKeys = {'creator': '💂 **criador**', 'producer': '🔴 **responsável**', 'creationDate': '📅 **dia de criação**', 
                    'modDate': '🕰️ **dia de modificação**', 'title': '#️⃣  **título**', 'author': '📕 **autor**', 'format': '⏹️ **formato**',
                    'subject': '🖊️ **assunto**', 'keywords': '#️⃣  **palavras-chave**', 'encryption': '🔑 **criptografia**'}
        keys = [key for key in list(dictKeys.keys())]
        for k, key in enumerate(keys):
            valueKey = dictKeys[key]
            metaKey = pdfMeta[key]
            if metaKey is None:
                metaKey = trace
            else:
                if len(metaKey.strip()) == 0:
                    metaKey = trace
            if k in [2, 3]:
                metaKey = configDate(metaKey)                
            st.markdown(f'{dictKeys[key]}: {metaKey}')
    config()
    
def configDate(datePdf):
    try:
        dateSplit = datePdf.split(':')
        dateStr = dateSplit[1][:14]
        year = dateStr[:4]
        month = dateStr[4:6]
        day = dateStr[6:8]
        hour = dateStr[8:10]
        minute = dateStr[10:12]
        second = dateStr[12:]
        dateStr = f'{day}/{month}/{year}, {hour}h{minute}min{second}s'
        return dateStr
    except:
        return datePdf
    
def exibeQrCode():
    @st.dialog('Dados')
    def config():
        nameUser = st.text_input(label='Nome', key=qrCodeKeys[0], placeholder=valuesReserve[0], 
                                      value='')
        phoneUser = st.text_input(label='Telefone', key=qrCodeKeys[1], placeholder=valuesReserve[1], value=''), 
        emailUser = st.text_input(label='Email', key=qrCodeKeys[2], placeholder=valuesReserve[2], value='')
        buttReturn = st.button('retornar')
        if buttReturn:
            for key in qrCodeKeys:
                del st.session_state[key]
            st.session_state[qrCodeKeys[0]] = nameUser
            st.session_state[qrCodeKeys[1]] = phoneUser
            st.session_state[qrCodeKeys[2]] = emailUser
            st.rerun()
    config()
                
@st.dialog(' ')
def config(str):
    st.text(str)  
    
@st.dialog('Critérios')
def windowAdd(numPgOne, numPgTwo):
    selModel = st.selectbox(label=f'Páginas a considerar no intervalo de :blue[**{numPgOne}**] a :blue[**{numPgTwo}**]', 
                            options=optionsSel, index=st.session_state[listKeys[5]])
    if selModel:
        del st.session_state[listKeys[5]]
        st.session_state[listKeys[5]] = optionsSel.index(selModel)
    if st.button('retornar'):
        st.rerun()
       
def iniFinally(mode):
    if mode == 0:
        for key in listKeys:
            if key not in st.session_state:
                try:
                    st.session_state[key] = dictKeys[key]
                except:
                    pass        
    else:
        try:
            for key in listKeys:
                del st.session_state[key]
        except:
            pass  
        iniFinally(0)
        st.rerun()
        
def seqPages(numPgOne, numPgTwo):
    valNum = st.session_state[listKeys[5]] 
    listPgs = [pg for pg in range(numPgOne, numPgTwo)]
    match valNum:
        case 1:
            listPgs = [pg for pg in range(numPgOne, numPgTwo) if (pg+1)%2==0]
        case 2:
            listPgs = [pg for pg in range(numPgOne, numPgTwo) if (pg+1)%2==1]
        case 3:
            listPgs = [pg for pg in range(numPgOne, numPgTwo)]
        case 4:
            listPgs = [pg for pg in range(numPgOne, numPgTwo) if (pg+1)%3==0]
        case 5:
            listPgs = [pg for pg in range(numPgOne, numPgTwo) if (pg+1)%4==0]
        case 6:
           listPgs = [pg for pg in range(numPgOne, numPgTwo) if (pg+1)%5==0]
        case 7:
           listPgs = [pg for pg in range(numPgOne, numPgTwo) if (pg+1)%10==0]
        case 8:
           listPgs = [pg for pg in range(numPgOne, numPgTwo) if (pg+1)%15==0] 
        case 9:
            listPgs = [pg for pg in range(numPgOne, numPgTwo) if (pg+1)%20==0]
    return listPgs            

def main():
    global uploadPdf
    global valMx
    global sufix
    sufix = ['']
    with st.container(border=6):
        uploadPdf = st.file_uploader('Selecionar arquivos PDF', 
                                     type=['pdf'], 
                                     accept_multiple_files=False)
        if uploadPdf is not None:
            pdfName = uploadPdf.name
            docPdf = pymupdf.open(stream=uploadPdf.read(), filetype="pdf")
            valMx = docPdf.page_count 
            valMxSize = round(uploadPdf.size/(1024**2), 2)
            if valMxSize < dictKeys[listKeys[3]]:
                dictKeys[listKeys[3]] = valMxSize
            colPgs, colPgOne, colPgTwo, colSlider, colSize, colMark, colPerson = st.columns([0.4, 1.35, 1.35, 2.3, 1.6, 2.7, 0.4], 
                                                                                vertical_alignment='bottom')
            buttToPages = colPgs.button(label=dictButts[keysButts[-2]][0], use_container_width=True, 
                                        icon=dictButts[keysButts[-2]][1], key=keysButts[-2], 
                                        help=dictButts[keysButts[-2]][-1])
            numPgOne = colPgOne.number_input(label='Página inicial  (:red[**1**])', key=listKeys[0], 
                                             min_value=1, max_value=valMx)
            numPgTwo = colPgTwo.number_input(label=f'Página final  (:red[**{valMx}**])', key=listKeys[1], 
                                             min_value=1, max_value=valMx)
            valPgAngle = colSlider.select_slider(label='Ângulo de rotação', options=valAngles, 
                                                 key=listKeys[2])    
            valPgSize = colSize.number_input(label='Tamanho para divisão (:red[**MB**])', key=listKeys[3], 
                                             min_value=dictKeys[listKeys[3]], step=dictKeys[listKeys[3]],  
                                             max_value=valMxSize)
            valPgMark = colMark.text_input(label="Marca d'água", key=listKeys[4], max_chars=50, 
                                           value=dictKeys[listKeys[4]], placeholder=nameApp)
            buttPerson = colPerson.button(label=dictButts[keysButts[-1]][0], use_container_width=True, 
                                          icon=dictButts[keysButts[-1]][1], key=keysButts[-1], 
                                          help=dictButts[keysButts[-1]][-1]) 
            colButtAct, colButtTxt, colButtSel, colButtDel, colButtClear = st.columns(5)
            buttPgAct = colButtAct.button(label=dictButts[keysButts[0]][0], key=keysButts[0], 
                                          use_container_width=True, icon=dictButts[keysButts[0]][1], 
                                          help=dictButts[keysButts[0]][-1])
            buttPgTxt = colButtTxt.button(label=dictButts[keysButts[1]][0], key=keysButts[1], 
                                          use_container_width=True, icon=dictButts[keysButts[1]][1], 
                                          help=dictButts[keysButts[1]][-1])
            buttPgSel = colButtSel.button(label=dictButts[keysButts[2]][0], key=keysButts[2], 
                                          use_container_width=True, icon=dictButts[keysButts[2]][1], 
                                          help=dictButts[keysButts[2]][-1])
            buttPgDel = colButtDel.button(label=dictButts[keysButts[3]][0], key=keysButts[3], 
                                          use_container_width=True, icon=dictButts[keysButts[3]][1], 
                                          help=dictButts[keysButts[3]][-1])
            buttPgClear = colButtClear.button(label=dictButts[keysButts[4]][0], key=keysButts[4], 
                                              use_container_width=True, icon=dictButts[keysButts[4]][1], 
                                              help=dictButts[keysButts[4]][-1])
            colButtUrl, colButtImg, colButtSize, colButtMark, colButtInfo = st.columns(5)
            buttPdfUrl = colButtUrl.button(label=dictButts[keysButts[5]][0], key=keysButts[5], 
                                           use_container_width=True, icon=dictButts[keysButts[5]][1], 
                                           help=dictButts[keysButts[5]][-1])
            buttPdfImg = colButtImg.button(label=dictButts[keysButts[6]][0], key=keysButts[6], 
                                           use_container_width=True, icon=dictButts[keysButts[6]][1], 
                                           help=dictButts[keysButts[6]][-1])
            buttPdfSize = colButtSize.button(label=dictButts[keysButts[7]][0], key=keysButts[7], 
                                             use_container_width=True, icon=dictButts[keysButts[7]][1], 
                                             help=dictButts[keysButts[7]][-1])
            buttPdfMark = colButtMark.button(label=dictButts[keysButts[8]][0], key=keysButts[8], 
                                             use_container_width=True, icon=dictButts[keysButts[8]][1], 
                                             help=dictButts[keysButts[8]][-1])
            buttPdfInfo =  colButtInfo.button(label=dictButts[keysButts[9]][0], key=keysButts[9], 
                                              use_container_width=True, icon=dictButts[keysButts[9]][1], 
                                              help=dictButts[keysButts[9]][-1])
            colTxtTable, colToTable, colToImg, colToPower, colCode = st.columns(5)
            buttTxtTable = colTxtTable.button(label=dictButts[keysButts[10]][0], key=keysButts[10], 
                                              use_container_width=True, icon=dictButts[keysButts[10]][1], 
                                              help=dictButts[keysButts[10]][-1])
            buttToWord = colToTable.button(label=dictButts[keysButts[11]][0], key=keysButts[11], 
                                           use_container_width=True, icon=dictButts[keysButts[11]][1], 
                                           help=dictButts[keysButts[11]][-1])
            buttToImg = colToImg.button(label=dictButts[keysButts[12]][0], key=keysButts[12], 
                                        use_container_width=True, icon=dictButts[keysButts[12]][1], 
                                        help=dictButts[keysButts[12]][-1])
            buttToPower = colToPower.button(label=dictButts[keysButts[13]][0], key=keysButts[13], 
                                            use_container_width=True, icon=dictButts[keysButts[13]][1], 
                                            help=dictButts[keysButts[13]][-1])   
            buttQrcode =  colCode.button(label=dictButts[keysButts[14]][0], key=keysButts[14], 
                                         use_container_width=True, icon=dictButts[keysButts[14]][1], 
                                         help=dictButts[keysButts[14]][-1])                                   
            if numPgTwo >= numPgOne: 
                numPgIni = numPgOne
                numPgFinal = numPgTwo
            else:
                numPgIni = numPgTwo
                numPgFinal = numPgOne 
            indexAng = valAngles.index(valPgAngle)
            exprPre = f'o intervalo de páginas {numPgOne} a {numPgTwo}.'
            if buttToPages:
                windowAdd(numPgOne, numPgTwo)
            if buttPgAct:  
                try:
                    expr = f'{dictButts[keysButts[0]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        extractPgs(docPdf, numPgIni, numPgFinal, 0, pdfName, indexAng)
                except:
                    config(f'😢 Divisão fracassada!\n🔴 arquivo {namePdf}, intervalo de páginas {numPgOne}-{numPgTwo}!') 
            if buttPerson:
                exibeQrCode()
            if buttPgTxt: 
                try:
                    expr = f'{dictButts[keysButts[1]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        selTxtUrlPgs(docPdf, numPgOne, numPgTwo, pdfName, 0, indexAng)
                except:
                     config(f'😢 Extração de texto fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttPgSel:
                try:
                    expr = f'{dictButts[keysButts[2]][2]} {pdfName} {exprPre}'
                    with st.spinner(expr):
                        selDelPgs(docPdf, numPgOne, numPgTwo, pdfName, 0, indexAng)
                except:
                    config(f'😢 Seleção de páginas fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttPgDel: 
                try:
                    expr = f'{dictButts[keysButts[3]][2]} {pdfName} {exprPre}'
                    with st.spinner(expr):
                        selDelPgs(docPdf, numPgOne, numPgTwo, pdfName, 1, indexAng)
                except:
                    config(f'😢 Deleção de páginas fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttPgClear: 
                del st.session_state[listKeys[5]]
                st.session_state[listKeys[5]] = 0
                iniFinally(1) 
            if buttPdfUrl:
                try:
                    expr = f'{dictButts[keysButts[5]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        sufix[0] = 'urls'
                        selTxtUrlPgs(docPdf, numPgOne, numPgTwo, pdfName, 1, indexAng)
                except:
                     config(f'😢 Extração de link fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttPdfImg: 
                try:     
                    expr = f'{dictButts[keysButts[6]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        sufix[0] = 'imgs'
                        selImgUrlsPgs(docPdf, numPgOne, numPgTwo, pdfName, 2, indexAng)
                except:
                    config(f'😢 Extração de imagens fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!') 
            if buttPdfSize:
                try:
                    expr = f'{dictButts[keysButts[7]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        selPgsSize(docPdf, numPgOne, numPgTwo, pdfName, indexAng, valPgSize)
                except:
                    config(f'😢 Divisão em pedaços fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttPdfMark:
                try:
                    if valPgMark.strip() == '':
                        valPgMark = nameApp 
                    expr = f'{dictButts[keysButts[8]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        selPdfMark(docPdf, numPgOne, numPgTwo, pdfName, indexAng, valPgMark)
                except:
                    config(f'😢 Marcação de páginas fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttPdfInfo:
                try:
                    exibeInfo(docPdf)
                except:
                    config(f'😢 Exibição fracassada!\n🔴 arquivo {pdfName}!')
            if buttPgClear: 
                del st.session_state[listKeys[5]]
                st.session_state[listKeys[5]] = 0
                iniFinally(1) 
            if buttTxtTable:
                try:
                    expr = f'{dictButts[keysButts[10]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        selTablesPgs(docPdf, numPgOne, numPgTwo, pdfName, indexAng)          
                except:
                    config(f'😢 Extração de tabelas fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttToWord:
                try:
                    expr = f'{dictButts[keysButts[11]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        selPdfToDocx(docPdf, numPgOne, numPgTwo, pdfName, indexAng)    
                except:
                    config(f'😢 Conversão de PDF em docx fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttToImg:
                try:
                    expr = f'{dictButts[keysButts[12]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        selPdfToImg(docPdf, numPgOne, numPgTwo, pdfName, indexAng)
                except: 
                    config(f'😢 Conversão de PDF em imagem fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttToPower:
                try:
                    expr = f'{dictButts[keysButts[13]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        selPdfToPPtx(docPdf, numPgOne, numPgTwo, pdfName, indexAng)                      
                except:
                    config(f'😢 Conversão de PDF em Power Point fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')
            if buttQrcode:
                try:
                    expr = f'{dictButts[keysButts[14]][2]} {pdfName} n{exprPre}'
                    with st.spinner(expr):
                        selPdfToQrcode(docPdf, numPgOne, numPgTwo, pdfName, indexAng)                        
                except:
                    config(f'😢 Inserção de QRcode fracassada!\n🔴 arquivo {pdfName}, intervalo de páginas {numPgOne}-{numPgTwo}!')    
                        
if __name__ == '__main__':
    global dictKeys, listKeys 
    global valAngles, valComps
    global countPg, optionsSel
    global namesTeste, nameApp 
    global qrCodeKeys, valuesReserve
    global dictButts, keysButts
    nameApp = 'Ferramentas/PDF'
    valAngles = ['-360°', '-270°', '-180°', '-90°', '0°', '90°', '180°', '270°', '360°']
    optionsSel = ['', 'pares', 'não pares', 'todos', 'de 3 em 3', 'de 4 em 4', 'de 5 em 5', 'de 10 em 10', 'de 15 em 15', 
                  'de 20 em 20'] 
    dictKeys = {'pgOne': 1, 
                'pgTwo': 1, 
                'pgAngle': valAngles[0], 
                'pgSize': 0.01, 
                'pgMark': '', 
                'selModelExtra': 0}
    listKeys = list(dictKeys.keys())
    dictButts = {'buttActIni': ['Corte/páginas', ':material/cut:', 'Dividindo o arquivo ', 'Secciona o arquivo de acordo com o intervalo de páginas.'], 
                 'buttTxt': ['Texto', ':material/description:', 'Extraindo texto do arquivo ', 'Extrai texto do arquivo e grava o resultado como txt.'],
                 'buttSel': ['Seleção', ':material/description:', 'Selecionando do arquivo ', 'Cria novo arquivo pdf com as páginas selecionadas.'], 
                 'buttDel': ['Deleção', ':material/delete:', 'Deletando do arquivo ', 'Deleta as páginas selecionadas.'], 
                 'buttClear': ['Limpeza', ':material/square:', 'Limpando os campos da tela.', 'Limpa os campos da tela, exceto o arquivo escolhido.'], 
                 'buttUrls': ['URLs', ':material/link:', 'Extraindo links/URLs do arquivo ', 'Pesquisa as URLs existentes no arquivo.'], 
                 'buttImgs': ['Imagens', ':material/image:', 'Extraindo imagens do arquivo', 'Extrai imagens do arquivo do arquivo e grava-as individualmente.'], 
                 'buttSize': ['Corte/tamanho', ':material/docs:', 'Dividindo por tamanho o arquivo ', 'Secciona o arquivo de acordo com o tamanho escolhido.'], 
                 'buttMark': ['Marcação', ':material/approval:', 'Marcando o rodapé do arquivo ', 'Insere marca de água nop rodapé do arquivo.'], 
                 'buttInfo': ['Informações', ':material/info:', 'Coligindo informações sobre o arquivo inteiro.', 'Exibe informações sobre o arquivo inteiro.'], 
                 'buttTxtTab': ['Texto/tabela', ':material/table:', 'Extraindo tabelas do arquivo ', 'Extrai tabelas existentes no treho selecionado.'], 
                 'buttToWord': ['Docx', ':material/transform:', 'Convertendo em Word o arquivo ', 'Converte em formato docx as páginas selecionadas do arquivo.'], 
                 'buttToImg': ['Imagem', ':material/modeling:', 'Convertendo em imagem (png) o arquivo ', 'Converte em formato jpg as páginas selecionadas.'], 
                 'buttToPower': ['Pptx', ':material/cycle:', 'Convertendo em slide do PowerPoint o arquivo ', 'Converte em slide do PowerPoint as páginas selecionadas.'], 
                 'buttQrcode': ['Qrcode', ':material/qr_code_2:', 'Inserindo qrcode no canto inferior direito do arquivo ', 'Insere qrcode no rodapé das páginas selecionadas.'], 
                 'buttPgs': ['', ':material/settings:', 'Exibindo opções de seleção de páginas do arquivo ', 'Exibe opções especiais de seleção de páginas.'],
                 'buttToPerson': ['', ':material/person_edit:', 'Abrindo campos a preencher para inserção do qrcode', 'Abre opções para preenchimento do qrcode.']}
    keysButts = list(dictButts.keys())
                
    countPg = []
    namesTeste = []
    dirBin = r'C:\Users\ACER\Documents\bin'
    valuesReserve = ['xxxxxxxx xxxxxxx', '+55xxxxxxxxxxx', 'xxxxxxxx@xxxxx.xxxx']
    qrCodeKeys = ['one', 'two', 'three']
    for key in qrCodeKeys:
        if key not in st.session_state:
            st.session_state[key] = ''    
    st.set_page_config(page_title=nameApp,  page_icon=":material/files:", 
                       layout='center')
    st.cache_data.clear() 
    iniFinally(0)
    with open('configuration.css') as f:
        css = f.read()
    st.markdown(f'<style>{css}</style>', unsafe_allow_html=True) 
    main()
